import google.generativeai as genai
from openpyxl import Workbook
from pptx import Presentation

"""
This code uses Google’s generativeai package to create a project plan based on a given project charter.
The plan is then processed into a structured list and saved as an Excel spreadsheet and a PowerPoint presentation.
The presentations slides use the generated tasks as titles and fill in other task details in the content.
"""


def generate(api_key, project_charter):
  
    genai.configure(api_key=api_key)
    prompt = f"""Act as a senior project manager with experience in software products. Write a Project Plan based on the Project Charter below:
    {project_charter}
    --
    END OF PROJECT CHARTER
    The output of the Project Plan needs to be in tabular format. The columns are:
    1) Task Name: Description of the task.
    2) Duration: Time required to complete the task.
    3) Dependencies: Other tasks that must be completed before this task.
    4) Status: Current status of the task (e.g., Not Started, In Progress, Completed).
    5) Resources: Tools, team members, software, infrastructure, etc., required for the task.
    Example of the output:
    Task Name|Duration|Dependencies|Status|Resources
    Creating Plan|1 day|None|Not Started|Leadership team
    Documentation|3 days|Creating Plan|Not Started|Development team, internal software
    Don't add any additional content or notes after you finish listing the tasks. Add at least 10 tasks and no extra content after. Now please write the Project Plan in tabular format:"""
  
    response = genai.generate_text(prompt=prompt)
    rows = response.result.split('\n')
    spreadsheet = []
  
    for row in rows:
        split_row = row.split('|')
        spreadsheet.append(split_row)
    print(spreadsheet)
        
    excel_save(spreadsheet)
    powerpoint_save(spreadsheet)
  
    return response.result


def excel_save(spreadsheet):
  
    wb = Workbook()
    ws = wb.active
    for row in spreadsheet:
        ws.append(row)
    wb.save('project_plan.xlsx')


def powerpoint_save(spreadsheet):
  
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1] # title and content layout
    spreadsheet = spreadsheet[4:]
  
    for row in spreadsheet:
        if len(row) < 5:
            continue
        slide = presentation.slides.add_slide(slide_layout)
        # title
        slide.placeholders[0].text = row[0]
        # content
        slide.placeholders[1].text = "Duration: " + row[1] + "\nDependencies: " + row[2] + "\nStatus: " + row[3] + "\nResources: " + row[4] 
    presentation.save('project_plan.pptx')
